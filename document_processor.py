#document_processor.py

import os
from typing import List, Dict

def extract_text_from_file(file_path: str) -> List[Dict]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        return _from_docx(file_path)
    elif ext == ".pptx":
        return _from_pptx(file_path)
    elif ext in (".txt", ".md"):
        return _from_text(file_path)
    elif ext == ".epub":
        return _from_epub(file_path)
    return []

def _from_docx(file_path):
    from docx import Document
    doc = Document(file_path)
    text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"text": text, "page_number": 1, "file_name": os.path.basename(file_path), "file_path": file_path, "total_pages": 1}]

def _from_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
        if text.strip():
            pages.append({"text": text, "page_number": i + 1, "file_name": os.path.basename(file_path), "file_path": file_path, "total_pages": len(prs.slides)})
    return pages

def _from_text(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [{"text": text, "page_number": 1, "file_name": os.path.basename(file_path), "file_path": file_path, "total_pages": 1}]

def _from_epub(file_path):
    import ebooklib
    from ebooklib import epub
    from html.parser import HTMLParser
    
    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []
        def handle_data(self, data):
            self.text.append(data)
    
    book = epub.read_epub(file_path)
    pages = []
    for i, item in enumerate(book.get_items_of_type(ebooklib.ITEM_DOCUMENT)):
        parser = TextExtractor()
        parser.feed(item.get_content().decode("utf-8", errors="ignore"))
        text = " ".join(parser.text).strip()
        if text:
            pages.append({"text": text, "page_number": i + 1, "file_name": os.path.basename(file_path), "file_path": file_path, "total_pages": i + 1})
    return pages